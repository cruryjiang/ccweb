"""
PPT 生成器模块
提供交互式 PPT 制作功能，支持模板管理和 URL 参考文档
"""
import os
import json
import re
import asyncio
import aiohttp
from dataclasses import dataclass, field, asdict
from typing import List, Dict, Optional, AsyncGenerator
from datetime import datetime
import tempfile
import shutil
from html.parser import HTMLParser
from urllib.parse import quote, unquote

# PPT 生成
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 文档解析
try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False


@dataclass
class SlideContent:
    """单页幻灯片内容"""
    title: str = ""
    subtitle: str = ""
    bullets: List[str] = field(default_factory=list)
    notes: str = ""  # 演讲者备注
    layout_type: str = "content"  # title, content, title_only, blank, etc.

    def to_dict(self) -> dict:
        return asdict(self)

    @classmethod
    def from_dict(cls, data: dict) -> "SlideContent":
        return cls(**data)


@dataclass
class PPTOutline:
    """PPT 大纲结构"""
    title: str = ""
    subtitle: str = ""
    slides: List[SlideContent] = field(default_factory=list)
    theme: str = "default"  # 主题风格
    source_doc: str = ""  # 源文档内容摘要

    def to_dict(self) -> dict:
        return {
            "title": self.title,
            "subtitle": self.subtitle,
            "theme": self.theme,
            "source_doc": self.source_doc,
            "slides": [s.to_dict() for s in self.slides]
        }

    @classmethod
    def from_dict(cls, data: dict) -> "PPTOutline":
        slides = [SlideContent.from_dict(s) for s in data.get("slides", [])]
        return cls(
            title=data.get("title", ""),
            subtitle=data.get("subtitle", ""),
            theme=data.get("theme", "default"),
            source_doc=data.get("source_doc", ""),
            slides=slides
        )


@dataclass
class PPTSession:
    """PPT 制作会话"""
    session_id: str
    outline: PPTOutline = field(default_factory=PPTOutline)
    current_slide_index: int = 0
    status: str = "idle"  # idle, outline_generated, editing, finalizing, completed
    created_at: float = field(default_factory=lambda: datetime.now().timestamp())
    source_file: str = ""  # 上传的源文档路径
    output_file: str = ""  # 生成的 PPT 文件路径
    template_file: str = ""  # 用户选择的模板文件名

    def to_dict(self) -> dict:
        return {
            "session_id": self.session_id,
            "outline": self.outline.to_dict(),
            "current_slide_index": self.current_slide_index,
            "status": self.status,
            "created_at": self.created_at,
            "source_file": self.source_file,
            "output_file": self.output_file,
            "template_file": self.template_file
        }


# 内存存储（生产环境应使用 Redis/数据库）
ppt_sessions: Dict[str, PPTSession] = {}

# PPT 文件存储目录（使用数据目录，避免写入只读的打包目录）
_data_dir = os.environ.get('CLAUDE_WEB_DATA_DIR', os.path.dirname(os.path.abspath(__file__)))
PPT_STORAGE_DIR = os.path.join(_data_dir, "data", "ppt_files")
os.makedirs(PPT_STORAGE_DIR, exist_ok=True)

# 模板存储目录
PPT_TEMPLATE_DIR = os.path.join(_data_dir, "data", "ppt_templates")
os.makedirs(PPT_TEMPLATE_DIR, exist_ok=True)


# ===================== HTML 文本提取器 =====================

class _HTMLTextExtractor(HTMLParser):
    """从 HTML 中提取纯文本"""
    def __init__(self):
        super().__init__()
        self._texts = []
        self._skip = False

    def handle_starttag(self, tag, attrs):
        if tag in ('script', 'style', 'nav', 'footer', 'header'):
            self._skip = True

    def handle_endtag(self, tag):
        if tag in ('script', 'style', 'nav', 'footer', 'header'):
            self._skip = False
        if tag in ('p', 'div', 'br', 'li', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'tr'):
            self._texts.append('\n')

    def handle_data(self, data):
        if not self._skip:
            text = data.strip()
            if text:
                self._texts.append(text)

    def get_text(self):
        return ' '.join(self._texts)


# ===================== 模板管理 =====================

def list_templates() -> List[Dict]:
    """列出所有可用的 PPT 模板"""
    templates = []
    if not os.path.exists(PPT_TEMPLATE_DIR):
        return templates

    for filename in sorted(os.listdir(PPT_TEMPLATE_DIR)):
        if filename.lower().endswith('.pptx') and not filename.startswith('.'):
            filepath = os.path.join(PPT_TEMPLATE_DIR, filename)
            stat = os.stat(filepath)
            name = os.path.splitext(filename)[0]
            templates.append({
                "name": name,
                "filename": filename,
                "size": stat.st_size,
                "modified": datetime.fromtimestamp(stat.st_mtime).isoformat(),
            })
    return templates


def save_template(filename: str, content: bytes) -> Dict:
    """保存上传的模板文件"""
    if not filename.lower().endswith('.pptx'):
        raise ValueError("模板文件必须是 .pptx 格式")

    # 清理文件名
    safe_name = re.sub(r'[^\w\-. \u4e00-\u9fff]', '_', filename)
    filepath = os.path.join(PPT_TEMPLATE_DIR, safe_name)

    with open(filepath, 'wb') as f:
        f.write(content)

    return {
        "name": os.path.splitext(safe_name)[0],
        "filename": safe_name,
        "size": len(content),
    }


def delete_template(filename: str) -> bool:
    """删除模板文件"""
    filepath = os.path.join(PPT_TEMPLATE_DIR, filename)
    if os.path.exists(filepath):
        os.remove(filepath)
        return True
    return False


def get_template_path(filename: str) -> Optional[str]:
    """获取模板文件的完整路径"""
    if not filename:
        return None
    filepath = os.path.join(PPT_TEMPLATE_DIR, filename)
    if os.path.exists(filepath):
        return filepath
    return None


# ===================== 在线模板搜索与下载 =====================

async def search_and_download_template(description: str) -> Dict:
    """
    从互联网搜索并下载免费 PPTX 模板。
    使用 DuckDuckGo 搜索，然后尝试从搜索结果页面中找到 .pptx 下载链接。
    """
    if not description or not description.strip():
        return {"status": "error", "message": "请输入模板描述"}

    timeout = aiohttp.ClientTimeout(total=90)
    headers = {
        'User-Agent': 'Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) '
                       'AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36',
        'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8',
        'Accept-Language': 'zh-CN,zh;q=0.9,en;q=0.8',
    }

    queries = [
        f"{description} free pptx template download",
        f"{description} powerpoint 模板 免费下载 pptx",
    ]

    pptx_urls = []
    page_urls = []

    async with aiohttp.ClientSession(timeout=timeout, headers=headers) as session:
        # Step 1: 搜索引擎查找模板链接
        for query in queries:
            try:
                search_url = f"https://html.duckduckgo.com/html/?q={quote(query)}"
                async with session.get(search_url) as resp:
                    if resp.status != 200:
                        continue
                    html = await resp.text()

                    # 提取直接 .pptx 链接
                    direct_links = re.findall(
                        r'https?://[^\s"\'<>]+\.pptx(?:\?[^\s"\'<>]*)?', html
                    )
                    pptx_urls.extend(direct_links)

                    # 提取页面链接（排除搜索引擎自身）
                    hrefs = re.findall(r'href="(/lite\?[^"]*|https?://[^"]+)"', html)
                    for href in hrefs:
                        if 'duckduckgo' in href or 'google' in href:
                            continue
                        if href.startswith('http'):
                            page_urls.append(href)

                if pptx_urls:
                    break
            except Exception:
                continue

        # Step 2: 访问搜索结果页面，提取 .pptx 下载链接
        if not pptx_urls and page_urls:
            for page_url in page_urls[:8]:
                try:
                    async with session.get(page_url, allow_redirects=True) as resp:
                        if resp.status != 200:
                            continue
                        ct = resp.headers.get('Content-Type', '')
                        # 如果直接就是 pptx 文件
                        if 'application/vnd.openxmlformats' in ct or page_url.endswith('.pptx'):
                            content = await resp.read()
                            if len(content) > 5000 and content[:2] == b'PK':
                                pptx_urls.insert(0, page_url)
                                # 直接保存已下载的内容
                                return await _save_downloaded_template(
                                    content, description, page_url
                                )
                            continue

                        page_html = await resp.text()
                        found = re.findall(
                            r'https?://[^\s"\'<>]+\.pptx(?:\?[^\s"\'<>]*)?', page_html
                        )
                        pptx_urls.extend(found)

                        # 也查找 download 链接
                        dl_links = re.findall(
                            r'href="(https?://[^"]*(?:download|\.pptx)[^"]*)"', page_html
                        )
                        for dl in dl_links:
                            if dl not in pptx_urls:
                                pptx_urls.append(dl)

                    if pptx_urls:
                        break
                except Exception:
                    continue

        # Step 3: 尝试下载找到的 .pptx 文件
        # 去重
        seen = set()
        unique_urls = []
        for u in pptx_urls:
            if u not in seen:
                seen.add(u)
                unique_urls.append(u)

        for url in unique_urls[:5]:
            try:
                async with session.get(url, allow_redirects=True) as resp:
                    if resp.status != 200:
                        continue
                    content = await resp.read()
                    if len(content) < 5000:
                        continue
                    if content[:2] != b'PK':
                        continue

                    return await _save_downloaded_template(content, description, url)
            except Exception:
                continue

    return {"status": "error", "message": "未找到合适的免费模板，请尝试其他描述或手动上传"}


async def _save_downloaded_template(
    content: bytes, description: str, source_url: str
) -> Dict:
    """保存下载的模板文件"""
    # 尝试验证是否为有效的 PPTX（可被 python-pptx 打开）
    tmp_path = os.path.join(PPT_TEMPLATE_DIR, f".tmp_{int(datetime.now().timestamp())}.pptx")
    try:
        with open(tmp_path, 'wb') as f:
            f.write(content)
        # 验证文件
        prs = Presentation(tmp_path)
        slide_count = len(prs.slides)
        layout_count = len(prs.slide_layouts)
    except Exception as e:
        if os.path.exists(tmp_path):
            os.remove(tmp_path)
        return {"status": "error", "message": f"下载的文件不是有效的 PPTX 模板: {e}"}

    # 生成安全文件名
    safe_name = re.sub(r'[^\w\- \u4e00-\u9fff]', '', description)[:30].strip()
    if not safe_name:
        safe_name = 'online_template'
    filename = f"{safe_name}.pptx"
    # 如果文件已存在，添加时间戳
    if os.path.exists(os.path.join(PPT_TEMPLATE_DIR, filename)):
        filename = f"{safe_name}_{int(datetime.now().timestamp())}.pptx"

    final_path = os.path.join(PPT_TEMPLATE_DIR, filename)
    os.rename(tmp_path, final_path)

    return {
        "status": "ok",
        "template": {
            "name": safe_name,
            "filename": filename,
            "size": len(content),
            "slide_count": slide_count,
            "layout_count": layout_count,
        },
        "message": f"模板下载成功：{safe_name}（{slide_count} 页，{layout_count} 种布局）",
    }


# ===================== URL 内容获取 =====================

async def fetch_url_content(url: str, max_length: int = 10000) -> str:
    """
    从 URL 获取文本内容作为参考文档
    支持 HTML 网页，自动提取正文文本
    """
    try:
        timeout = aiohttp.ClientTimeout(total=30)
        async with aiohttp.ClientSession(timeout=timeout) as session:
            async with session.get(url, headers={
                'User-Agent': 'Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36'
            }) as resp:
                if resp.status != 200:
                    return f"[无法获取 URL 内容，状态码: {resp.status}]"

                content_type = resp.headers.get('Content-Type', '')
                text = await resp.text(errors='replace')

                if 'text/html' in content_type or '<html' in text[:500].lower():
                    extractor = _HTMLTextExtractor()
                    extractor.feed(text)
                    text = extractor.get_text()

                # 截断过长内容
                if len(text) > max_length:
                    text = text[:max_length] + "\n...[内容已截断]"
                return text
    except Exception as e:
        return f"[获取 URL 内容失败: {str(e)}]"


def extract_text_from_document(file_path: str) -> str:
    """
    从文档中提取文本内容
    支持 TXT, MD, DOCX, PDF
    """
    ext = os.path.splitext(file_path)[1].lower()

    if ext in ['.txt', '.md', '.markdown']:
        with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
            return f.read()

    elif ext == '.docx':
        if not DOCX_AVAILABLE:
            raise ImportError("python-docx 未安装，无法读取 Word 文档")
        doc = DocxDocument(file_path)
        texts = []
        for para in doc.paragraphs:
            if para.text.strip():
                texts.append(para.text)
        return '\n'.join(texts)

    elif ext == '.pdf':
        if not PDF_AVAILABLE:
            raise ImportError("PyPDF2 未安装，无法读取 PDF 文档")
        text = []
        with open(file_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                text.append(page.extract_text() or "")
        return '\n'.join(text)

    else:
        raise ValueError(f"不支持的文件格式: {ext}")


def generate_outline_prompt(source_text: str, user_requirement: str = "", reference_text: str = "") -> str:
    """
    生成用于创建 PPT 大纲的 AI Prompt

    Args:
        source_text: 上传的源文档内容
        user_requirement: 用户需求描述
        reference_text: 从 URL 获取的参考文档内容
    """
    prompt = f"""请根据以下内容生成一个 PPT 大纲。

要求：
1. 设计一个清晰的标题和副标题
2. 规划 5-15 页幻灯片（根据内容复杂度决定）
3. 每页包含：标题、要点（3-5 个 bullet points）、页面类型建议
4. 确保逻辑流畅，覆盖核心内容
5. 使用 JSON 格式输出

"""
    if user_requirement:
        prompt += f"用户需求：{user_requirement}\n\n"

    if source_text:
        prompt += f"源文档内容：\n{source_text[:8000]}\n\n"

    if reference_text:
        prompt += f"参考文档内容：\n{reference_text[:8000]}\n\n"

    prompt += """请按以下 JSON 格式输出：
{
  "title": "PPT 主标题",
  "subtitle": "副标题",
  "theme": "business|academic|creative|minimal",
  "slides": [
    {
      "title": "页面标题",
      "subtitle": "副标题（可选，用于标题页）",
      "bullets": ["要点1", "要点2", "要点3"],
      "layout_type": "title|content|title_only|two_content"
    }
  ]
}

注意：
- 第一页通常是 title 类型（标题页）
- 内容页使用 content 类型
- layout_type 可选值：title, content, title_only, two_content, blank
"""
    return prompt


def generate_slide_content_prompt(outline: PPTOutline, slide_index: int, source_text: str = "") -> str:
    """
    生成用于完善单页幻灯片内容的 AI Prompt
    """
    if slide_index < 0 or slide_index >= len(outline.slides):
        return ""

    slide = outline.slides[slide_index]

    prompt = f"""请完善以下 PPT 页面的内容：

页面标题：{slide.title}
当前要点：
"""
    for bullet in slide.bullets:
        prompt += f"- {bullet}\n"

    prompt += f"\n页面类型：{slide.layout_type}\n"

    if source_text:
        prompt += f"\n参考原文（相关部分）：\n{source_text[:3000]}\n"

    prompt += """
请完善内容：
1. 优化标题，使其更吸引人
2. 扩展要点，使其更详细、有价值（每个要点 1-2 句话）
3. 添加演讲者备注（notes），帮助演讲者理解这页的重点
4. 如果需要，建议这页是否需要配图或图表

请按以下 JSON 格式输出：
{
  "title": "优化后的标题",
  "subtitle": "副标题（如有）",
  "bullets": ["详细要点1", "详细要点2", "详细要点3"],
  "notes": "演讲者备注",
  "suggestions": "设计建议，如配图建议"
}
"""
    return prompt


def _extract_json_object(text: str) -> Optional[dict]:
    """
    从文本中提取包含 "slides" 键的 JSON 对象。
    使用多种策略，按可靠性从高到低尝试。
    """
    # 策略 1: 从 markdown 代码块中提取
    code_block_match = re.search(r'```(?:json)?\s*(\{[\s\S]*?\})\s*```', text)
    if code_block_match:
        try:
            data = json.loads(code_block_match.group(1))
            if isinstance(data, dict) and "slides" in data:
                return data
        except json.JSONDecodeError:
            pass

    # 策略 2: 使用平衡括号匹配，找到包含 "slides" 的最外层 JSON 对象
    # 从每个 { 开始尝试配对
    for match in re.finditer(r'\{', text):
        start = match.start()
        depth = 0
        i = start
        while i < len(text):
            ch = text[i]
            if ch == '{':
                depth += 1
            elif ch == '}':
                depth -= 1
                if depth == 0:
                    candidate = text[start:i + 1]
                    try:
                        data = json.loads(candidate)
                        if isinstance(data, dict) and "slides" in data:
                            return data
                    except json.JSONDecodeError:
                        pass
                    break
            elif ch == '"':
                # 跳过字符串内容
                i += 1
                while i < len(text) and text[i] != '"':
                    if text[i] == '\\':
                        i += 1  # 跳过转义字符
                    i += 1
            i += 1

    # 策略 3: 回退到贪婪匹配（原始方式）
    json_match = re.search(r'\{[\s\S]*\}', text)
    if json_match:
        try:
            data = json.loads(json_match.group())
            if isinstance(data, dict):
                return data
        except json.JSONDecodeError:
            pass

    return None


def parse_outline_from_ai_response(response: str) -> Optional[PPTOutline]:
    """
    从 AI 响应中解析 PPT 大纲
    """
    try:
        data = _extract_json_object(response)
        if not data:
            print(f"解析大纲失败: 未在响应中找到有效 JSON（响应长度={len(response)}）")
            return None

        # 构建 PPTOutline
        slides = []
        for slide_data in data.get("slides", []):
            slide = SlideContent(
                title=slide_data.get("title", ""),
                subtitle=slide_data.get("subtitle", ""),
                bullets=slide_data.get("bullets", []),
                notes=slide_data.get("notes", ""),
                layout_type=slide_data.get("layout_type", "content")
            )
            slides.append(slide)

        if not slides:
            print(f"解析大纲失败: slides 为空")
            return None

        return PPTOutline(
            title=data.get("title", ""),
            subtitle=data.get("subtitle", ""),
            theme=data.get("theme", "default"),
            source_doc=data.get("source_doc", ""),
            slides=slides
        )
    except Exception as e:
        print(f"解析大纲失败: {e}")
        return None


def parse_slide_content_from_ai_response(response: str) -> Optional[SlideContent]:
    """
    从 AI 响应中解析单页幻灯片内容
    """
    try:
        data = _extract_json_object(response)
        if not data:
            # 回退：slide 数据可能不含 "slides" 键，用简单提取
            code_block = re.search(r'```(?:json)?\s*(\{[\s\S]*?\})\s*```', response)
            if code_block:
                try:
                    data = json.loads(code_block.group(1))
                except json.JSONDecodeError:
                    pass
            if not data:
                json_match = re.search(r'\{[^{}]*"title"[^{}]*\}', response)
                if json_match:
                    try:
                        data = json.loads(json_match.group())
                    except json.JSONDecodeError:
                        pass
            if not data:
                return None

        return SlideContent(
            title=data.get("title", ""),
            subtitle=data.get("subtitle", ""),
            bullets=data.get("bullets", []),
            notes=data.get("notes", ""),
            layout_type=data.get("layout_type", "content")
        )
    except Exception as e:
        print(f"解析幻灯片内容失败: {e}")
        return None


def get_or_create_ppt_session(session_id: str) -> PPTSession:
    """获取或创建 PPT 会话"""
    if session_id not in ppt_sessions:
        ppt_sessions[session_id] = PPTSession(session_id=session_id)
    return ppt_sessions[session_id]


def delete_ppt_session(session_id: str):
    """删除 PPT 会话及其文件"""
    if session_id in ppt_sessions:
        session = ppt_sessions[session_id]

        # 删除源文件
        if session.source_file and os.path.exists(session.source_file):
            try:
                os.remove(session.source_file)
            except:
                pass

        # 删除输出文件
        if session.output_file and os.path.exists(session.output_file):
            try:
                os.remove(session.output_file)
            except:
                pass

        del ppt_sessions[session_id]


def cleanup_old_sessions(max_age_hours: int = 24):
    """清理过期的会话和文件"""
    current_time = datetime.now().timestamp()
    to_delete = []

    for session_id, session in ppt_sessions.items():
        age_hours = (current_time - session.created_at) / 3600
        if age_hours > max_age_hours:
            to_delete.append(session_id)

    for session_id in to_delete:
        delete_ppt_session(session_id)


def _find_layout_by_type(slide_layouts, layout_type: str):
    """
    根据布局类型在模板的 slide_layouts 中查找最佳匹配的布局。
    优先按名称匹配，其次按占位符结构匹配，最后回退到安全索引。
    """
    num_layouts = len(slide_layouts)

    # 名称关键词匹配映射（中英文 + 常见变体）
    name_keywords = {
        "title": ["title slide", "标题幻灯片", "封面", "section divider", "divider"],
        "content": ["title and content", "标题和内容", "内容", "1 placeholder", "content"],
        "title_only": ["title only", "仅标题", "blank"],
        "two_content": ["two content", "比较", "双栏", "2 placeholder", "2 column"],
        "blank": ["blank", "空白"],
    }

    # 按名称匹配
    keywords = name_keywords.get(layout_type, [])
    for layout in slide_layouts:
        layout_name = layout.name.lower()
        for kw in keywords:
            if kw in layout_name:
                return layout

    # 按占位符结构匹配
    if layout_type == "content":
        # 找一个有 OBJECT 类型或 BODY 类型占位符的布局（内容页）
        for layout in slide_layouts:
            has_title = False
            has_content = False
            for ph in layout.placeholders:
                if ph.placeholder_format.idx == 0:
                    has_title = True
                if ph.placeholder_format.type in (2, 7):  # BODY=2, OBJECT=7
                    has_content = True
            if has_title and has_content:
                return layout

    # 按索引回退（使用安全索引，防止越界）
    index_map = {
        "title": 0,
        "content": min(1, num_layouts - 1),
        "title_only": min(2, num_layouts - 1),
        "two_content": min(3, num_layouts - 1),
        "blank": num_layouts - 1,
    }
    idx = index_map.get(layout_type, min(1, num_layouts - 1))
    return slide_layouts[max(0, min(idx, num_layouts - 1))]


def _get_placeholder_by_idx(slide, idx):
    """通过 idx 获取占位符（兼容所有 python-pptx 版本）"""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def _find_subtitle_placeholder(slide):
    """在幻灯片中查找副标题占位符"""
    for ph in slide.placeholders:
        name_lower = ph.name.lower()
        # 按名称查找
        if "sub-title" in name_lower or "subtitle" in name_lower or "副标题" in name_lower:
            return ph
        # 按类型查找（SUBTITLE = 4）
        if ph.placeholder_format.type == 4:
            return ph
    # 回退到 idx=10 或 idx=11（常见自定义模板），最后 idx=1
    for idx in [10, 11, 1]:
        ph = _get_placeholder_by_idx(slide, idx)
        if ph and ph.placeholder_format.idx != 0:
            return ph
    return None


def _find_content_placeholder(slide):
    """在幻灯片中查找内容占位符（用于 bullet points）"""
    # 优先按名称查找包含 "Content" 的占位符
    for ph in slide.placeholders:
        name_lower = ph.name.lower()
        if "content" in name_lower and ph.placeholder_format.idx != 0:
            return ph
    # 其次查找 OBJECT 类型（type=7）
    for ph in slide.placeholders:
        if ph.placeholder_format.type == 7 and ph.placeholder_format.idx != 0:
            return ph
    # 再次查找 BODY 类型（type=2），排除 sub-title、footnote、source 等
    for ph in slide.placeholders:
        if ph.placeholder_format.type == 2 and ph.placeholder_format.idx != 0:
            name_lower = ph.name.lower()
            if not any(kw in name_lower for kw in ["sub-title", "subtitle", "footnote", "source", "sequence"]):
                return ph
    # 最后回退到 idx=1
    ph = _get_placeholder_by_idx(slide, 1)
    if ph:
        return ph
    return None


def _remove_all_slides(prs):
    """安全地删除 Presentation 中所有已有幻灯片，保留布局和主题"""
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        if rId:
            prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)


def _remove_slides_from(prs, keep_count: int):
    """删除 Presentation 中 keep_count 之后的所有幻灯片"""
    sldIdLst = prs.slides._sldIdLst
    all_sldIds = list(sldIdLst)
    for sldId in all_sldIds[keep_count:]:
        rId = sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        if rId:
            prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)


def _fill_cover_slide(slide, outline: 'PPTOutline'):
    """
    填充模板封面页：扫描所有 shapes，将含 "XX" 占位文本的 shape 替换为实际内容。
    - 含 "汇报" / "报告" 等关键词的大文本 → 替换 XX 为 outline.title
    - 含 "月" / "年" 的日期文本 → 替换为当前日期
    - 其他含 XX 的文本 → 替换为 outline.subtitle
    """
    now = datetime.now()
    date_str = f"{now.year}年{now.month:02d}月"

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        full_text = shape.text_frame.text
        if not full_text.strip():
            continue

        # 日期类文本：含 "月" 或 "年" 且有 "XX"（如 "2026年XX月"）
        if ("月" in full_text or "年" in full_text) and "XX" in full_text:
            month_str = f"{now.month:02d}"
            for para in shape.text_frame.paragraphs:
                _replace_xx_in_paragraph(para, month_str)
            continue

        # 标题类文本：含 "汇报" / "报告" / "总结" 等
        if "XX" in full_text and any(kw in full_text for kw in ["汇报", "报告", "总结", "分享", "介绍"]):
            for para in shape.text_frame.paragraphs:
                _replace_xx_in_paragraph(para, outline.title or "工作汇报")
            continue

        # 其他含 XX 的文本 → 用 subtitle 替换
        if "XX" in full_text:
            replacement = outline.subtitle or outline.title or ""
            for para in shape.text_frame.paragraphs:
                _replace_xx_in_paragraph(para, replacement)


def _replace_xx_in_paragraph(para, replacement: str):
    """替换段落中所有 run 里的 'XX' 为 replacement，保留原有格式"""
    for run in para.runs:
        if "XX" in run.text:
            run.text = run.text.replace("XX", replacement)


def _replace_text_in_paragraph(para, old_full_text: str, new_text: str):
    """
    替换段落的整体文本。
    保留第一个 run 的格式，清空其余 run。
    """
    if not para.runs:
        return
    para.runs[0].text = new_text
    for run in para.runs[1:]:
        run.text = ""


def _fill_toc_slide(slide, outline: 'PPTOutline'):
    """
    填充模板目录页：用 outline.slides 的标题生成目录内容。
    扫描 shapes，找到含有编号列表或 "目录" 的 shape 进行填充。
    """
    # 收集章节标题（跳过 title 类型的封面页）
    toc_items = []
    for i, s in enumerate(outline.slides):
        if s.layout_type == "title" and i == 0:
            continue
        toc_items.append(s.title)

    if not toc_items:
        return

    # 查找目录内容区域：排除仅含 "目录" 二字的标题 shape
    target_shape = None
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        # 跳过只有 "目录" 标题的 shape
        if text == "目录" or text == "CONTENTS" or text == "Contents":
            continue
        # 找有多行文本或编号的 shape（目录内容区）
        if len(shape.text_frame.paragraphs) > 1 or len(text) > 10:
            target_shape = shape
            break

    if target_shape is None:
        # 没找到合适的，用第一个非标题的有文本 shape
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if text and text != "目录" and text != "CONTENTS":
                target_shape = shape
                break

    if target_shape is None:
        return

    # 填充目录项
    tf = target_shape.text_frame
    # 保存第一个段落的格式作为模板
    template_para = tf.paragraphs[0] if tf.paragraphs else None
    template_font_size = None
    template_font_color = None
    template_font_bold = None
    if template_para and template_para.runs:
        run = template_para.runs[0]
        template_font_size = run.font.size
        template_font_color = run.font.color.rgb if run.font.color and run.font.color.type is not None else None
        template_font_bold = run.font.bold

    tf.clear()
    for i, item in enumerate(toc_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{i + 1}. {item}"
        p.space_after = Pt(6)
        # 应用模板格式
        if p.runs:
            run = p.runs[0]
            if template_font_size:
                run.font.size = template_font_size
            if template_font_color:
                run.font.color.rgb = template_font_color
            if template_font_bold is not None:
                run.font.bold = template_font_bold


def _add_content_slide(prs, slide_content, use_template: bool):
    """添加一页内容幻灯片并填充标题和要点"""
    # 根据内容类型选择最佳布局
    if slide_content.layout_type == "title_only":
        target_type = "title_only"
    elif slide_content.layout_type == "two_content":
        target_type = "two_content"
    else:
        target_type = "content"

    layout = _find_layout_by_type(prs.slide_layouts, target_type)
    slide = prs.slides.add_slide(layout)

    # 设置标题
    if slide.shapes.title:
        slide.shapes.title.text = slide_content.title

    # 内容页：设置 bullet points
    content_ph = _find_content_placeholder(slide)
    if content_ph and slide_content.bullets:
        tf = content_ph.text_frame
        tf.clear()
        for j, bullet in enumerate(slide_content.bullets):
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
    elif slide_content.bullets:
        # 无合适占位符，添加文本框作为回退
        left = Inches(1)
        top = Inches(2)
        width = prs.slide_width - Inches(2)
        height = prs.slide_height - Inches(3)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True
        for j, bullet in enumerate(slide_content.bullets):
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {bullet}"
            p.space_after = Pt(8)

    # 添加演讲者备注
    if slide_content.notes:
        try:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_content.notes
        except Exception:
            pass  # 模板的 notes 布局缺少 idx=1 占位符，跳过备注

    return slide


def create_ppt_file(outline: PPTOutline, output_path: str, template_path: str = None) -> str:
    """
    使用 python-pptx 创建 PPT 文件

    Args:
        outline: PPT 大纲
        output_path: 输出文件路径
        template_path: 模板文件路径（可选）

    返回: 输出文件的完整路径
    """
    use_template = template_path and os.path.exists(template_path)

    if use_template:
        prs = Presentation(template_path)
        template_slide_count = len(prs.slides)

        # 模板使用策略：保留封面/目录页，替换内容
        if template_slide_count >= 1:
            # 填充封面页
            _fill_cover_slide(list(prs.slides)[0], outline)

        if template_slide_count >= 2:
            # 填充目录页
            _fill_toc_slide(list(prs.slides)[1], outline)

        # 删除模板中 keep_count 之后的示例页
        keep_count = min(2, template_slide_count)
        _remove_slides_from(prs, keep_count)

        # 逐页添加内容幻灯片（跳过 outline 中的第一页 title 页，因为已用模板封面）
        for i, slide_content in enumerate(outline.slides):
            if i == 0 and slide_content.layout_type == "title":
                continue  # 封面页已由模板提供
            _add_content_slide(prs, slide_content, use_template=True)

    else:
        prs = Presentation()
        # 设置幻灯片尺寸为 16:9
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        for i, slide_content in enumerate(outline.slides):
            if slide_content.layout_type == "title" or i == 0:
                # 标题页
                layout = _find_layout_by_type(prs.slide_layouts, "title")
                slide = prs.slides.add_slide(layout)
                if slide.shapes.title:
                    slide.shapes.title.text = slide_content.title
                subtitle_ph = _find_subtitle_placeholder(slide)
                if subtitle_ph and slide_content.subtitle:
                    subtitle_ph.text_frame.text = slide_content.subtitle
                # 添加演讲者备注
                if slide_content.notes:
                    try:
                        notes_slide = slide.notes_slide
                        notes_slide.notes_text_frame.text = slide_content.notes
                    except Exception:
                        pass
            else:
                _add_content_slide(prs, slide_content, use_template=False)

    # 添加结束页
    if outline.slides and outline.slides[-1].title not in ["谢谢", "Thank You", "结束", "Q&A"]:
        ending_layout = _find_layout_by_type(prs.slide_layouts, "blank")
        ending_slide = prs.slides.add_slide(ending_layout)

        # 添加感谢文本
        left = Inches(2)
        top = Inches(3)
        width = Inches(9)
        height = Inches(1.5)

        textbox = ending_slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.text = "谢谢观看"
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(44)
        p.font.bold = True

    # 保存文件
    prs.save(output_path)
    return output_path


async def stream_ppt_creation_process(
    session_id: str,
    claude_stream_func,
    message: str = ""
) -> AsyncGenerator[str, None]:
    """
    流式获取 PPT 创建过程

    claude_stream_func: 调用 Claude CLI 的流式函数
    """
    session = get_or_create_ppt_session(session_id)

    if session.status == "idle":
        # 第一步：生成大纲
        yield json.dumps({
            "type": "status",
            "message": "正在分析文档并生成 PPT 大纲...",
            "step": "outline"
        })

        source_text = ""
        if session.source_file:
            try:
                source_text = extract_text_from_document(session.source_file)
            except Exception as e:
                yield json.dumps({
                    "type": "error",
                    "message": f"读取源文档失败: {str(e)}"
                })
                return

        prompt = generate_outline_prompt(source_text, message)

        # 调用 Claude 生成大纲
        response_text = ""
        async for chunk in claude_stream_func(prompt, session_id):
            try:
                data = json.loads(chunk[6:])  # 去掉 "data: " 前缀
                if data.get("type") == "content":
                    response_text += data.get("text", "")
                elif data.get("type") == "done":
                    break
            except:
                pass

        # 解析大纲
        outline = parse_outline_from_ai_response(response_text)
        if outline:
            session.outline = outline
            session.status = "outline_generated"
            yield json.dumps({
                "type": "outline_ready",
                "outline": outline.to_dict(),
                "message": f"大纲已生成，共 {len(outline.slides)} 页"
            })
        else:
            yield json.dumps({
                "type": "error",
                "message": "无法解析 AI 返回的大纲，请重试"
            })

    elif session.status == "outline_generated":
        # 处于大纲已生成状态，等待用户确认或修改
        yield json.dumps({
            "type": "status",
            "message": "请确认大纲或提出修改意见",
            "step": "review_outline",
            "outline": session.outline.to_dict()
        })

    elif session.status == "editing":
        # 编辑单页内容
        current_slide = session.outline.slides[session.current_slide_index]
        yield json.dumps({
            "type": "status",
            "message": f"正在编辑第 {session.current_slide_index + 1} 页...",
            "step": "editing_slide",
            "slide_index": session.current_slide_index,
            "slide": current_slide.to_dict()
        })

    elif session.status == "finalizing":
        # 生成最终 PPT 文件
        yield json.dumps({
            "type": "status",
            "message": "正在生成 PPT 文件...",
            "step": "generating_file"
        })

        output_filename = f"ppt_{session_id}_{int(datetime.now().timestamp())}.pptx"
        output_path = os.path.join(PPT_STORAGE_DIR, output_filename)

        try:
            create_ppt_file(session.outline, output_path)
            session.output_file = output_path
            session.status = "completed"

            yield json.dumps({
                "type": "ppt_ready",
                "message": "PPT 生成完成！",
                "filename": output_filename,
                "download_url": f"/api/ppt/download/{output_filename}",
                "slide_count": len(session.outline.slides)
            })
        except Exception as e:
            yield json.dumps({
                "type": "error",
                "message": f"生成 PPT 文件失败: {str(e)}"
            })


# PPT 命令帮助文本
PPT_COMMAND_HELP = """
📊 PPT 制作命令帮助

使用方法：
/ppt [需求描述]  - 进入 PPT 制作模式

工作流程：
1. 上传源文档（可选）- 支持 TXT, MD, DOCX, PDF
2. 输入需求描述      - 如："制作一个产品经理年终总结 PPT"
3. AI 生成大纲       - 预览并确认 PPT 结构
4. 逐页编辑内容      - 修改每页的标题和要点
5. 生成 PPT 文件     - 下载 .pptx 文件

编辑命令：
- 确认 / 下一步     - 确认当前步骤继续
- 修改第N页         - 编辑指定页面
- 添加一页          - 在最后添加新页面
- 删除第N页         - 删除指定页面
- 完成              - 生成最终 PPT 文件
- 退出              - 退出 PPT 制作模式
"""
